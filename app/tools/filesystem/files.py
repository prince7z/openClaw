from __future__ import annotations

import csv
import json
import zipfile
from pathlib import Path
from typing import Any, Literal
from xml.dom import minidom
from xml.etree import ElementTree as ET

import ebooklib
import yaml
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from langchain.tools import tool
from langchain_core.runnables import RunnableConfig
from markdownify import markdownify as html_to_markdown
from openpyxl import load_workbook
from app.tools.filesystem.supported import SUPPORTED_DOCUMENT_EXTENSIONS, UNSUPPORTED_BINARY_EXTENSIONS
import fitz
from pptx import Presentation

from app.tools.filesystem._common import (
    ensure_parent_directory,
    line_count,
    looks_binary,
    markdown_table,
    make_result,
    mime_type_for_path,
    safe_read_bytes,
    to_path,
    validate_paths,
    validate_text_content,
)



ODT_TEXT_NS = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"


def _code_block(language: str, content: str) -> str:
    """Wrap text in a fenced Markdown code block."""
    return f"```{language}\n{content}\n```"


def _normalize_markdown(content: str) -> str:
    """Normalize output for LLM consumption."""
    return content.replace("\r\n", "\n").replace("\r", "\n").strip()


def _document_payload(path: Path, content: str) -> dict[str, Any]:
    """Build the structured response payload for read_file()."""
    normalized = _normalize_markdown(content)
    return {
        "content": normalized,
        "format": "markdown",
        "extension": path.suffix.lower(),
        "line_count": line_count(normalized),
        "size": path.stat().st_size,
        "mime_type": mime_type_for_path(path),
    }


def _read_plain_text(path: Path, encoding: str) -> str:
    """Read a plain text or markdown file."""
    return path.read_text(encoding=encoding)


def _read_json(path: Path, encoding: str) -> str:
    """Read JSON as pretty Markdown code block."""
    parsed = json.loads(path.read_text(encoding=encoding))
    pretty = json.dumps(parsed, indent=2, ensure_ascii=False)
    return _code_block("json", pretty)


def _read_yaml(path: Path, encoding: str) -> str:
    """Read YAML as pretty Markdown code block."""
    loaded = yaml.safe_load(path.read_text(encoding=encoding))
    pretty = yaml.safe_dump(loaded, sort_keys=False, allow_unicode=True).strip()
    return _code_block("yaml", pretty)


def _read_xml(path: Path, encoding: str) -> str:
    """Read XML as pretty Markdown code block."""
    raw = path.read_text(encoding=encoding)
    pretty = minidom.parseString(raw.encode(encoding)).toprettyxml(indent="  ")
    return _code_block("xml", pretty.strip())


def _read_csv(path: Path, encoding: str) -> str:
    """Read CSV as a Markdown table."""
    with path.open("r", encoding=encoding, newline="") as handle:
        rows = list(csv.reader(handle))

    if not rows:
        return "(empty csv)"

    headers = rows[0]
    data_rows = rows[1:]
    table = markdown_table(data_rows, headers=headers)
    return table or "(empty csv)"


def _read_html(path: Path, encoding: str) -> str:
    """Read HTML and convert it to Markdown."""
    raw = path.read_text(encoding=encoding)
    soup = BeautifulSoup(raw, "html.parser")
    for node in soup(["script", "style"]):
        node.decompose()
    markdown = html_to_markdown(str(soup), heading_style="ATX")
    return markdown.strip()


def _read_pdf(path: Path) -> str:
    """Read PDF content and convert it to Markdown-like text."""
    document = fitz.open(str(path))
    try:
        sections: list[str] = []
        for page_number, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            if text:
                sections.append(f"## Page {page_number}\n\n{text}")
        return "\n\n".join(sections).strip()
    finally:
        document.close()


def _read_docx(path: Path) -> str:
    """Read DOCX content and convert it to Markdown."""
    document = Document(str(path))
    sections: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        style_name = getattr(getattr(paragraph, "style", None), "name", "") or ""
        style_lower = style_name.lower()
        if style_lower.startswith("heading"):
            digits = "".join(character for character in style_name if character.isdigit())
            level = max(1, min(6, int(digits) if digits else 1))
            sections.append(f"{'#' * level} {text}")
        elif "bullet" in style_lower or "list" in style_lower:
            sections.append(f"- {text}")
        else:
            sections.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        sections.append(f"## Table {table_index}")
        sections.append(markdown_table(rows[1:], headers=rows[0]))

    return "\n\n".join(section for section in sections if section.strip()).strip()


def _read_xlsx(path: Path) -> str:
    """Read XLSX content and convert each sheet into Markdown tables."""
    workbook = load_workbook(str(path), data_only=True)
    sections: list[str] = []

    try:
        for sheet in workbook.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                sections.append(f"## Sheet: {sheet.title}\n\n(empty sheet)")
                continue

            first_row = ["" if value is None else str(value) for value in rows[0]]
            if all(not cell for cell in first_row):
                headers = [f"Column {index + 1}" for index in range(len(first_row))]
                data_rows = [["" if value is None else value for value in row] for row in rows]
            else:
                headers = first_row
                data_rows = [["" if value is None else value for value in row] for row in rows[1:]]

            table = markdown_table(data_rows, headers=headers)
            section = f"## Sheet: {sheet.title}\n\n{table or '(empty sheet)'}"
            sections.append(section)

        return "\n\n".join(section for section in sections if section.strip()).strip()
    finally:
        workbook.close()


def _read_pptx(path: Path) -> str:
    """Read PPTX content and convert each slide into Markdown."""
    presentation = Presentation(str(path))
    sections: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = getattr(slide.shapes, "title", None)
        title = ""
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title = title_shape.text.strip()

        heading = f"## Slide {slide_number}" + (f": {title}" if title else "")
        slide_lines: list[str] = [heading]

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph_text = paragraph.text.strip()
                if not paragraph_text:
                    continue
                indent = "  " * getattr(paragraph, "level", 0)
                slide_lines.append(f"{indent}- {paragraph_text}")

        if len(slide_lines) == 1:
            slide_lines.append("(empty slide)")

        sections.append("\n".join(slide_lines))

    return "\n\n".join(sections).strip()


def _extract_odt_text(path: Path) -> str:
    """Read ODT content from content.xml and convert it to Markdown."""
    with zipfile.ZipFile(path) as archive:
        with archive.open("content.xml") as content_file:
            tree = ET.parse(content_file)

    sections: list[str] = []
    for element in tree.iter():
        if element.tag == f"{ODT_TEXT_NS}h":
            level_text = element.attrib.get(f"{ODT_TEXT_NS}outline-level", "1")
            try:
                level = max(1, min(6, int(level_text)))
            except ValueError:
                level = 1
            text = "".join(element.itertext()).strip()
            if text:
                sections.append(f"{'#' * level} {text}")
        elif element.tag == f"{ODT_TEXT_NS}p":
            text = "".join(element.itertext()).strip()
            if text:
                sections.append(text)

    return "\n\n".join(sections).strip()


def _read_epub(path: Path) -> str:
    """Read EPUB content and convert document items to Markdown."""
    book = epub.read_epub(str(path))
    sections: list[str] = []

    for item in book.get_items():
        if item.get_type() != ebooklib.ITEM_DOCUMENT:
            continue
        raw_html = item.get_content().decode("utf-8", errors="ignore")
        soup = BeautifulSoup(raw_html, "html.parser")
        for node in soup(["script", "style"]):
            node.decompose()
        markdown = html_to_markdown(str(soup), heading_style="ATX").strip()
        if markdown:
            sections.append(f"## {item.get_name()}\n\n{markdown}")

    return "\n\n".join(sections).strip()


def _read_generic_text(path: Path, encoding: str) -> str:
    """Read an unknown text file as Markdown/plain text."""
    return _read_plain_text(path, encoding)


def _extract_markdown_content(path: Path, encoding: str) -> str:
    """Dispatch to the appropriate extractor based on file extension."""
    extension = path.suffix.lower()
    if extension in {".txt", ".md", ".markdown"}:
        return _read_plain_text(path, encoding)
    if extension == ".csv":
        return _read_csv(path, encoding)
    if extension == ".json":
        return _read_json(path, encoding)
    if extension in {".yaml", ".yml"}:
        return _read_yaml(path, encoding)
    if extension in {".xml"}:
        return _read_xml(path, encoding)
    if extension in {".html", ".htm"}:
        return _read_html(path, encoding)
    if extension == ".pdf":
        return _read_pdf(path)
    if extension == ".docx":
        return _read_docx(path)
    if extension == ".xlsx":
        return _read_xlsx(path)
    if extension == ".pptx":
        return _read_pptx(path)
    if extension == ".odt":
        return _extract_odt_text(path)
    if extension == ".epub":
        return _read_epub(path)
    return _read_generic_text(path, encoding)


@tool
def read_file(
    path: str,
    encoding: str = "utf-8",
    max_size_mb: int = 50,
    config: RunnableConfig = None
) -> dict[str, Any]:
    """Read a file."""
    try:
        target = to_path(path, config=config)
        if not target.exists():
            return make_result(False, None, f"Path does not exist: {target}")
        if not target.is_file():
            return make_result(False, None, f"Path is not a file: {target}")

        if not isinstance(max_size_mb, int) or max_size_mb < 0:
            return make_result(False, None, "max_size_mb must be a non-negative integer")

        size_limit = max_size_mb * 1024 * 1024
        file_size = target.stat().st_size
        if file_size > size_limit:
            return make_result(False, None, f"File exceeds max_size_mb limit ({max_size_mb} MB)")

        extension = target.suffix.lower()
        if extension in UNSUPPORTED_BINARY_EXTENSIONS:
            return make_result(False, None, "Unsupported binary format")

        if extension not in SUPPORTED_DOCUMENT_EXTENSIONS:
            sample = safe_read_bytes(target)
            if looks_binary(sample):
                return make_result(False, None, "Unsupported binary format")

        content = _extract_markdown_content(target, encoding)
        payload = _document_payload(target, content)
        return make_result(True, payload, None)
    except UnicodeDecodeError:
        return make_result(False, None, "Unsupported binary format")
    except (zipfile.BadZipFile, ET.ParseError, json.JSONDecodeError, yaml.YAMLError, ValueError, OSError, RuntimeError) as exc:
        return make_result(False, None, str(exc))
    except Exception as exc:
        return make_result(False, None, str(exc))



@tool
def write_file(
    path: str,
    content: str,
    mode: Literal["overwrite", "append"] = "overwrite",
    encoding: str = "utf-8",
    config: RunnableConfig = None
) -> dict[str, Any]:
    """Write or append content to a file."""
    try:
        target = to_path(path, config=config)
        text = validate_text_content(content)

        if target.exists() and target.is_dir():
            return make_result(False, None, f"Path is a directory: {target}")

        ensure_parent_directory(target)
        if mode == "append":
            with target.open("a", encoding=encoding) as handle:
                handle.write(text)
        else:
            target.write_text(text, encoding=encoding)

        return make_result(True, str(target.resolve(strict=False)), None)
    except Exception as exc:
        return make_result(False, None, str(exc))



# Internal helper function, not exposed as tool
def create_file(path: str, exist_ok: bool = False, config: RunnableConfig = None, **kwargs: Any) -> dict[str, Any]:
    """Create an empty file and any necessary parent directories.

    Args:
        path: File path to create.
        exist_ok: Whether to allow existing files.

    Returns:
        The absolute file path that was created or already existed.
    """
    if kwargs:
        unexpected = ", ".join(f"'{k}'" for k in kwargs)
        return make_result(
            False,
            None,
            f"create_file got unexpected argument(s): {unexpected}. "
            f"To write content to a file, please use the write_file tool instead."
        )
    try:
        target = to_path(path, config=config)
        if target.exists():
            if target.is_dir():
                return make_result(False, None, f"Path is a directory: {target}")
            if not exist_ok:
                return make_result(False, None, f"File already exists: {target}")
            return make_result(True, str(target.resolve(strict=False)), None)

        ensure_parent_directory(target)
        target.touch(exist_ok=True)
        return make_result(True, str(target.resolve(strict=False)), None)
    except Exception as exc:
        return make_result(False, None, str(exc))

