from __future__ import annotations

import csv
import json
import os
import stat
import sys
import tempfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from app.tools.filesystem import (  # noqa: E402
    read_file,
    write_file,
    manage_file,
    list_files,
    search_files,
)

from docx import Document  # noqa: E402
from ebooklib import epub  # noqa: E402
from odf.opendocument import OpenDocumentText  # noqa: E402
from odf.text import H, P  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from pptx import Presentation  # noqa: E402
import fitz  # noqa: E402

REPORT_PATH = ROOT / "report.md"


@dataclass
class CheckResult:
    name: str
    status: str
    details: str


def invoke(tool: Any, **kwargs: Any) -> Any:
    """Invoke a LangChain tool or plain callable in a uniform way."""
    if hasattr(tool, "invoke"):
        return tool.invoke(kwargs)
    return tool(**kwargs)


def as_json_text(value: Any) -> str:
    """Serialize values for the markdown report."""
    return json.dumps(value, indent=2, ensure_ascii=False, default=str)


def expect(condition: bool, message: str) -> None:
    """Raise an assertion error when a check fails."""
    if not condition:
        raise AssertionError(message)


def record(name: str, status: str, details: str) -> CheckResult:
    """Create a report row."""
    return CheckResult(name=name, status=status, details=details)


def success(name: str, result: Any, validator: Callable[[dict[str, Any]], None] | None = None) -> CheckResult:
    """Validate a structured filesystem tool response."""
    if not isinstance(result, dict):
        return record(name, "FAIL", f"Expected dict result, got {type(result).__name__}: {result!r}")

    if not result.get("success"):
        return record(name, "FAIL", f"Tool reported failure: {result.get('error')}")

    if validator is not None:
        try:
            validator(result)
        except AssertionError as exc:
            return record(name, "FAIL", str(exc))

    return record(name, "PASS", as_json_text(result.get("data")))


def failure(name: str, result: Any, expected_fragment: str | None = None) -> CheckResult:
    """Validate an expected failure response."""
    if not isinstance(result, dict):
        return record(name, "FAIL", f"Expected dict result, got {type(result).__name__}: {result!r}")

    if result.get("success"):
        return record(name, "FAIL", f"Expected failure but tool succeeded: {as_json_text(result.get('data'))}")

    error_text = str(result.get("error"))
    if expected_fragment and expected_fragment.lower() not in error_text.lower():
        return record(name, "FAIL", f"Unexpected error text: {error_text}")

    return record(name, "PASS", error_text)


def skipped(name: str, reason: str) -> CheckResult:
    """Record a skipped check."""
    return record(name, "SKIP", reason)


def write_report(results: list[CheckResult], sandbox: Path) -> None:
    """Write a markdown report summarizing the smoke test run."""
    passed = sum(1 for result in results if result.status == "PASS")
    skipped_count = sum(1 for result in results if result.status == "SKIP")
    failed = sum(1 for result in results if result.status == "FAIL")
    generated_at = datetime.now(timezone.utc).isoformat()

    lines = [
        "# Filesystem Tool Regression Report",
        "",
        f"- Generated at: {generated_at}",
        f"- Workspace root: {ROOT}",
        f"- Sandbox used: {sandbox}",
        f"- Total checks: {len(results)}",
        f"- Passed: {passed}",
        f"- Skipped: {skipped_count}",
        f"- Failed: {failed}",
        "",
        "## Results",
        "",
        "| Check | Status | Details |",
        "| --- | --- | --- |",
    ]

    for result in results:
        details = result.details.replace("\n", " ").replace("|", "\\|")
        lines.append(f"| {result.name} | {result.status} | {details} |")

    lines.extend(
        [
            "",
            "## Summary",
            "",
            "This script exercises the filesystem tool package against an isolated temporary workspace.",
            "",
            "Skipped checks indicate a platform limitation such as symlink creation permissions.",
        ]
    )

    REPORT_PATH.write_text("\n".join(lines), encoding="utf-8")


def create_text(path: Path, content: str) -> Path:
    """Create a text file with UTF-8 content."""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


def create_csv_sample(path: Path) -> Path:
    """Create a CSV sample file."""
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["name", "value"])
        writer.writerow(["alpha", "1"])
        writer.writerow(["beta", "2"])
    return path


def create_json_sample(path: Path) -> Path:
    """Create a JSON sample file."""
    return create_text(path, json.dumps({"title": "Sample", "items": [1, 2, 3]}, ensure_ascii=False))


def create_yaml_sample(path: Path) -> Path:
    """Create a YAML sample file."""
    import yaml

    return create_text(path, yaml.safe_dump({"title": "Sample", "items": [1, 2, 3]}, sort_keys=False, allow_unicode=True))


def create_xml_sample(path: Path) -> Path:
    """Create an XML sample file."""
    xml = """<?xml version="1.0" encoding="UTF-8"?>\n<report><title>Sample</title><item>One</item></report>"""
    return create_text(path, xml)


def create_html_sample(path: Path) -> Path:
    """Create an HTML sample file."""
    html = """<html><head><title>Sample</title></head><body><h1>Heading</h1><p>Hello <strong>world</strong>.</p></body></html>"""
    return create_text(path, html)


def create_pdf_sample(path: Path) -> Path:
    """Create a PDF sample file with visible text."""
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF Sample\nHello from PDF")
    document.save(str(path))
    document.close()
    return path


def create_docx_sample(path: Path) -> Path:
    """Create a DOCX sample file."""
    document = Document()
    document.add_heading("DOCX Sample", level=1)
    document.add_paragraph("Hello from DOCX")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Answer"
    table.cell(1, 1).text = "42"
    document.save(str(path))
    return path


def create_xlsx_sample(path: Path) -> Path:
    """Create an XLSX sample file."""
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Key", "Value"])
    sheet.append(["Answer", 42])
    workbook.save(str(path))
    workbook.close()
    return path


def create_pptx_sample(path: Path) -> Path:
    """Create a PPTX sample file."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(72, 72, 600, 200)
    frame = textbox.text_frame
    frame.text = "PPTX Sample"
    paragraph = frame.add_paragraph()
    paragraph.text = "Hello from PPTX"
    presentation.save(str(path))
    return path


def create_odt_sample(path: Path) -> Path:
    """Create an ODT sample file."""
    document = OpenDocumentText()
    document.text.addElement(H(outlinelevel=1, text="ODT Sample"))
    document.text.addElement(P(text="Hello from ODT"))
    document.save(str(path))
    return path


def create_epub_sample(path: Path) -> Path:
    """Create an EPUB sample file."""
    book = epub.EpubBook()
    book.set_identifier("sample-book")
    book.set_title("Sample EPUB")
    book.set_language("en")

    chapter = epub.EpubHtml(title="Chapter 1", file_name="chap_01.xhtml", lang="en")
    chapter.content = "<html><body><h1>Chapter 1</h1><p>Hello from EPUB</p></body></html>"
    book.add_item(chapter)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", chapter]
    epub.write_epub(str(path), book)
    return path


def create_binary_sample(path: Path, magic: bytes) -> Path:
    """Create a small binary file."""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(magic + b"\x00\x01\x02\x03")
    return path


def make_sparse_large_file(path: Path, size_mb: int = 101) -> Path:
    """Create a sparse file larger than the default read limit."""
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("wb") as handle:
        handle.seek(size_mb * 1024 * 1024)
        handle.write(b"x")
    return path


def make_read_only(path: Path) -> None:
    """Make a file read-only across platforms as much as possible."""
    path.chmod(stat.S_IREAD | stat.S_IRGRP | stat.S_IROTH)


def restore_write(path: Path) -> None:
    """Restore write permissions to a file."""
    path.chmod(stat.S_IWRITE | stat.S_IREAD)


def build_sample_workspace(sandbox: Path) -> dict[str, Path]:
    """Create the sample files needed for the regression run."""
    samples: dict[str, Path] = {}

    docs = sandbox / "docs"
    edge = sandbox / "edge"
    bulk = sandbox / "bulk"
    empty = sandbox / "empty"
    ops = sandbox / "ops"

    manage_file.invoke({"action": "create", "path": str(docs), "kind": "directory"})
    manage_file.invoke({"action": "create", "path": str(edge), "kind": "directory"})
    manage_file.invoke({"action": "create", "path": str(bulk), "kind": "directory"})
    manage_file.invoke({"action": "create", "path": str(empty), "kind": "directory"})
    manage_file.invoke({"action": "create", "path": str(ops), "kind": "directory"})

    samples["empty"] = empty
    samples["docs"] = docs
    samples["edge"] = edge
    samples["bulk"] = bulk
    samples["ops"] = ops

    samples["txt"] = create_text(docs / "sample.txt", "Hello text file")
    samples["md"] = create_text(docs / "sample.md", "# Markdown\n\nPlain markdown text.")
    samples["csv"] = create_csv_sample(docs / "sample.csv")
    samples["json"] = create_json_sample(docs / "sample.json")
    samples["xml"] = create_xml_sample(docs / "sample.xml")
    samples["yaml"] = create_yaml_sample(docs / "sample.yaml")
    samples["yml"] = create_yaml_sample(docs / "sample.yml")
    samples["html"] = create_html_sample(docs / "sample.html")
    samples["pdf"] = create_pdf_sample(docs / "sample.pdf")
    samples["docx"] = create_docx_sample(docs / "sample.docx")
    samples["xlsx"] = create_xlsx_sample(docs / "sample.xlsx")
    samples["pptx"] = create_pptx_sample(docs / "sample.pptx")
    samples["odt"] = create_odt_sample(docs / "sample.odt")
    samples["epub"] = create_epub_sample(docs / "sample.epub")

    samples["hidden"] = create_text(docs / ".hidden.txt", "hidden content")
    samples["unicode"] = create_text(docs / "unicodé_文件.txt", "unicode content")
    samples["readonly"] = create_text(edge / "readonly.txt", "read only content")
    samples["binary_exe"] = create_binary_sample(edge / "program.exe", b"MZ")
    samples["binary_png"] = create_binary_sample(edge / "image.png", b"\x89PNG\r\n\x1a\n")
    samples["large"] = make_sparse_large_file(edge / "huge.txt")

    for index in range(1100):
        create_text(bulk / f"file_{index:04d}.txt", f"file {index}")

    target = edge / "symlink_target.txt"
    create_text(target, "symlink target")
    samples["symlink_target"] = target
    samples["symlink"] = edge / "symlink.txt"
    samples["copy_source"] = create_text(ops / "copy_source.txt", "copy source")
    samples["copy_dest"] = ops / "copy_dest.txt"
    samples["move_source"] = create_text(ops / "move_source.txt", "move source")
    samples["move_dest"] = ops / "move_dest.txt"
    samples["trash_me"] = create_text(ops / "trash_me.txt", "trash me")
    samples["delete_dir"] = ops / "delete_me"

    return samples


def run_tool_checks(samples: dict[str, Path]) -> list[CheckResult]:
    """Execute the regression checks."""
    results: list[CheckResult] = []

    def add(name: str, result: Any, validator: Callable[[dict[str, Any]], None] | None = None) -> None:
        results.append(success(name, result, validator))

    def add_failure(name: str, result: Any, fragment: str | None = None) -> None:
        results.append(failure(name, result, fragment))

    add("empty_directory_list", list_files.invoke({"path": str(samples["empty"]), "view": "list"}), lambda result: expect(result["data"] == [], "empty directory should list no entries"))
    add("empty_directory_tree", list_files.invoke({"path": str(samples["empty"]), "view": "tree"}), lambda result: expect(result["data"]["children"] == [], "empty directory tree should have no children"))

    add("hidden_file_excluded", list_files.invoke({"path": str(samples["docs"]), "view": "list"}), lambda result: expect(all(item["name"] != ".hidden.txt" for item in result["data"]), "hidden file should be excluded by default"))

    add("unicode_filename_read", read_file.invoke({"path": str(samples["unicode"])}), lambda result: expect(result["data"]["content"] == "unicode content", "unicode file content mismatch"))

    add("bulk_directory_count", list_files.invoke({"path": str(samples["bulk"]), "recursive": False}), lambda result: expect(len(result["data"]) >= 1000, "bulk directory should contain 1000+ files"))
    add("bulk_directory_search_name", search_files.invoke({"path": str(samples["bulk"]), "query": "file_10", "search_type": "name"}), lambda result: expect(any("file_10" in item for item in result["data"]), "search should find bulk files by name"))
    add("bulk_directory_search_glob", search_files.invoke({"path": str(samples["bulk"]), "query": "file_1*.txt"}), lambda result: expect(any(Path(item).name.startswith("file_1") for item in result["data"]), "glob should match file_1* files"))

    symlink_path = samples["symlink"]
    try:
        if symlink_path.exists() or symlink_path.is_symlink():
            symlink_path.unlink()
        symlink_path.symlink_to(samples["symlink_target"])
        add("symlink_read_file", read_file.invoke({"path": str(symlink_path)}), lambda result: expect(result["data"]["content"] == "symlink target", "read_file should resolve symlink targets"))
    except (OSError, NotImplementedError) as exc:
        results.append(skipped("symlink_read_file", f"Symlink creation not available: {exc}"))

    add("read_txt", read_file.invoke({"path": str(samples["txt"])}), lambda result: expect(result["data"]["content"] == "Hello text file", "txt extraction mismatch"))
    add("read_md", read_file.invoke({"path": str(samples["md"])}), lambda result: expect(result["data"]["content"].startswith("# Markdown"), "markdown should be returned as-is"))
    add("read_csv", read_file.invoke({"path": str(samples["csv"])}), lambda result: expect("| name | value |" in result["data"]["content"], "csv should be converted to a markdown table"))
    add("read_json", read_file.invoke({"path": str(samples["json"])}), lambda result: expect(result["data"]["content"].startswith("```json"), "json should be pretty formatted as a code block"))
    add("read_xml", read_file.invoke({"path": str(samples["xml"])}), lambda result: expect(result["data"]["content"].startswith("```xml"), "xml should be returned as a markdown code block"))
    add("read_yaml", read_file.invoke({"path": str(samples["yaml"])}), lambda result: expect(result["data"]["content"].startswith("```yaml"), "yaml should be formatted as a code block"))
    add("read_yml", read_file.invoke({"path": str(samples["yml"])}), lambda result: expect(result["data"]["content"].startswith("```yaml"), "yml should be formatted as a code block"))
    add("read_html", read_file.invoke({"path": str(samples["html"])}), lambda result: expect("Heading" in result["data"]["content"], "html should be converted to markdown text"))
    add("read_pdf", read_file.invoke({"path": str(samples["pdf"])}), lambda result: expect("PDF Sample" in result["data"]["content"], "pdf should extract text"))
    add("read_docx", read_file.invoke({"path": str(samples["docx"])}), lambda result: expect("DOCX Sample" in result["data"]["content"], "docx should extract heading text"))
    add("read_xlsx", read_file.invoke({"path": str(samples["xlsx"])}), lambda result: expect("Sheet: Data" in result["data"]["content"], "xlsx should render a sheet markdown section"))
    add("read_pptx", read_file.invoke({"path": str(samples["pptx"])}), lambda result: expect("PPTX Sample" in result["data"]["content"], "pptx should extract slide text"))
    add("read_odt", read_file.invoke({"path": str(samples["odt"])}), lambda result: expect("ODT Sample" in result["data"]["content"], "odt should extract text"))
    add("read_epub", read_file.invoke({"path": str(samples["epub"])}), lambda result: expect("Chapter 1" in result["data"]["content"] or "Hello from EPUB" in result["data"]["content"], "epub should extract chapter content"))

    add_failure("read_binary_exe", read_file.invoke({"path": str(samples["binary_exe"])}), "Unsupported binary format")
    add_failure("read_binary_png", read_file.invoke({"path": str(samples["binary_png"])}), "Unsupported binary format")
    add_failure("large_file_default_limit", read_file.invoke({"path": str(samples["large"])}), "max_size_mb")

    add("write_file_append", write_file.invoke({"path": str(samples["txt"]), "content": "\nExtra line", "mode": "append"}), lambda result: expect(Path(result["data"]).exists(), "append should succeed"))

    readonly_file = samples["readonly"]
    make_read_only(readonly_file)
    try:
        add_failure("write_read_only_file", write_file.invoke({"path": str(readonly_file), "content": "changed", "mode": "overwrite"}), None)
    finally:
        restore_write(readonly_file)

    delete_dir = samples["ops"] / "delete_me"
    manage_file.invoke({"action": "create", "path": str(delete_dir), "kind": "directory"})
    create_text(delete_dir / "child.txt", "child")
    add("delete_directory_recursive", manage_file.invoke({"action": "delete", "path": str(delete_dir)}), lambda result: expect(not Path(result["data"]).exists(), "recursive delete should remove the directory"))

    add("create_directory_kind", manage_file.invoke({"action": "create", "path": str(samples["ops"] / "created_dir"), "kind": "directory"}), lambda result: expect(Path(result["data"]).exists(), "manage_file create directory should succeed"))
    add("create_file_kind", manage_file.invoke({"action": "create", "path": str(samples["ops"] / "created_file"), "kind": "file"}), lambda result: expect(Path(result["data"]).exists(), "manage_file create file should succeed"))

    add("copy_regular", manage_file.invoke({"action": "copy", "path": str(samples["copy_source"]), "target": str(samples["ops"] / "copied.txt")}), lambda result: expect(Path(result["data"]).exists(), "copy should create the destination file"))
    add("move_regular", manage_file.invoke({"action": "move", "path": str(samples["move_source"]), "target": str(samples["ops"] / "moved.txt")}), lambda result: expect(Path(result["data"]).exists(), "move should create the destination file"))
    add("rename_regular", manage_file.invoke({"action": "rename", "path": str(samples["ops"] / "moved.txt"), "target": "renamed.txt"}), lambda result: expect(Path(result["data"]).name == "renamed.txt", "rename should change the base name"))

    add("search_content", search_files.invoke({"path": str(samples["docs"]), "query": "Markdown", "search_type": "content"}), lambda result: expect(any("sample.md" in item for item in result["data"]), "content search should find Markdown in sample.md"))

    return results


def run_suite() -> int:
    """Execute all filesystem regression checks and write the report."""
    with tempfile.TemporaryDirectory(prefix="openclaw-filesystem-") as temp_dir:
        sandbox = Path(temp_dir).resolve()
        samples = build_sample_workspace(sandbox)
        results = run_tool_checks(samples)
        write_report(results, sandbox)

    return 1 if any(result.status == "FAIL" for result in results) else 0


if __name__ == "__main__":
    sys.exit(run_suite())
